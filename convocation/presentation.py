import os
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from typing import Optional, Union, List, Dict, Any

from convocation.config import CardLayoutConfig
from convocation.utils import find_photo
from convocation.card import generate_student_card

def create_presentation_from_cards(
    student_data: Union[pd.DataFrame, List[Dict[str, Any]]],
    template_path: str,
    image_folder: str,
    output_folder: str,
    output_pptx: str,
    layout_config: CardLayoutConfig,
    logo_path: Optional[str] = None,
    slide_width_inch: float = 20.0,
    slide_height_inch: float = 11.25,
    draw_batch: bool = True,
    overwrite_cards: bool = False
) -> None:
    """
    Renders student cards and packages them into a PowerPoint slide presentation.
    If a card image already exists in output_folder, it will reuse it unless overwrite_cards=True.
    """
    # Initialize PPTX presentation
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    prs.slide_width = Inches(slide_width_inch)
    prs.slide_height = Inches(slide_height_inch)
    
    os.makedirs(output_folder, exist_ok=True)
    template_base = Image.open(template_path)
    
    logo_image = None
    # Resolve logo image
    logo_file = logo_path or layout_config.logo_path
    if logo_file and os.path.exists(logo_file):
        logo_image = Image.open(logo_file).convert("RGB")
        
    rows = student_data.to_dict(orient="records") if isinstance(student_data, pd.DataFrame) else student_data
    
    for row in rows:
        roll_no = str(row.get('Roll No', '')).strip()
        if not roll_no:
            continue
            
        output_img_path = os.path.join(output_folder, f"{roll_no}.png")
        
        # Check if we can reuse the existing image
        if os.path.exists(output_img_path) and not overwrite_cards:
            print(f"Card for {roll_no} already exists, reusing...")
        else:
            photo = find_photo(roll_no, image_folder)
            if not photo:
                print(f"Photo not found for Roll No: {roll_no}, skipping slide generation...")
                continue
                
            # Generate the card image
            card = generate_student_card(
                row=row,
                template_base=template_base,
                layout_config=layout_config,
                photo_image=photo,
                logo_image=logo_image,
                draw_batch=draw_batch
            )
            card.save(output_img_path)
            print(f"Generated card - {roll_no}")
            
        # Add slide containing the card
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            output_img_path, 
            Inches(0), Inches(0), 
            width=Inches(slide_width_inch), 
            height=Inches(slide_height_inch)
        )
        
    prs.save(output_pptx)
    print(f"Saved PowerPoint presentation to {output_pptx}")
